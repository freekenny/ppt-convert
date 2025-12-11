import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import io
import json
import re
import fitz  # PyMuPDF
from streamlit_paste_button import paste_image_button

# ================= 工具函数 =================

def hex_to_rgb(hex_color):
    try:
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    except:
        return RGBColor(0, 0, 0)

def extract_images(uploaded_files, pasted_image):
    """提取图片逻辑：支持文件上传和剪贴板"""
    all_images = []
    # 1. 优先处理粘贴的图片
    if pasted_image: 
        all_images.append(pasted_image)
    
    # 2. 处理上传的文件
    if uploaded_files:
        for f in uploaded_files:
            if f.type == "application/pdf":
                try:
                    doc = fitz.open(stream=f.read(), filetype="pdf")
                    for p in range(len(doc)):
                        # zoom=2 保证 PDF 转图片清晰度
                        pix = doc.load_page(p).get_pixmap(matrix=fitz.Matrix(2,2))
                        all_images.append(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))
                except: pass
            else:
                try: 
                    f.seek(0)
                    all_images.append(Image.open(f))
                except: pass
    return all_images

def analyze_layout(api_key, img, model_name):
    """Gemini 分析逻辑"""
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel(model_name)
    
    # Prompt: 强制合并复杂图形
    prompt = """
    Analyze this slide image for PowerPoint reconstruction.
    
    Task: Identify "text_block" and "visual_element".
    
    CRITICAL RULE FOR IMAGES:
    If there is a complex diagram, infographic, or overlapping illustration (e.g., a central timeline with icons, or a heart with waves), **GROUP THEM into ONE LARGE "visual_element"**. 
    DO NOT split a complex chart into 10 small icons. Capture the WHOLE context area.
    
    Return JSON array:
    [
        {
            "type": "text_block",
            "content": "Text",
            "box_2d": [ymin, xmin, ymax, xmax],
            "style": { "font_size": 12, "is_bold": false, "color_hex": "#000000", "alignment": "left" }
        },
        {
            "type": "visual_element",
            "box_2d": [ymin, xmin, ymax, xmax] 
        }
    ]
    """
    try:
        response = model.generate_content([prompt, img], generation_config={"response_mime_type": "application/json"})
        return json.loads(re.sub(r'^```json|```$', '', response.text.strip()))
    except: return []

def create_pptx(image_list, api_key, model_name):
    """PPT 生成逻辑"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    prog = st.progress(0); status = st.empty()
    
    for idx, img in enumerate(image_list):
        status.text(f"Processing slide {idx+1}/{len(image_list)}...")
        layout = analyze_layout(api_key, img, model_name)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        w, h = img.size
        
        # 分离图层
        visuals = [x for x in layout if x['type'] == 'visual_element']
        texts = [x for x in layout if x['type'] == 'text_block']
        
        # 1. 画图片 (底层) + Padding 缓冲
        for item in visuals:
            box = item.get('box_2d', [0,0,0,0])
            ymin, xmin, ymax, xmax = box
            
            # 向四周扩充 1.5% 防止切坏边缘
            pad_x = (xmax - xmin) * 0.015 
            pad_y = (ymax - ymin) * 0.015
            
            crop_x1 = max(0, (xmin/1000 * w) - pad_x)
            crop_y1 = max(0, (ymin/1000 * h) - pad_y)
            crop_x2 = min(w, (xmax/1000 * w) + pad_x)
            crop_y2 = min(h, (ymax/1000 * h) + pad_y)
            
            ppt_l = (crop_x1 / w) * prs.slide_width
            ppt_t = (crop_y1 / h) * prs.slide_height
            ppt_w = ((crop_x2 - crop_x1) / w) * prs.slide_width
            ppt_h = ((crop_y2 - crop_y1) / h) * prs.slide_height

            if crop_x2 > crop_x1 and crop_y2 > crop_y1:
                try:
                    cropped = img.crop((crop_x1, crop_y1, crop_x2, crop_y2))
                    buf = io.BytesIO(); cropped.save(buf, format="PNG"); buf.seek(0)
                    slide.shapes.add_picture(buf, ppt_l, ppt_t, width=ppt_w, height=ppt_h)
                except: pass

        # 2. 写文字 (顶层)
        for item in texts:
            box = item.get('box_2d', [0,0,0,0])
            ymin, xmin, ymax, xmax = box
            
            ppt_l = (xmin/1000) * prs.slide_width
            ppt_t = (ymin/1000) * prs.slide_height
            ppt_w = ((xmax-xmin)/1000) * prs.slide_width
            ppt_h = ((ymax-ymin)/1000) * prs.slide_height
            
            tx = slide.shapes.add_textbox(ppt_l, ppt_t, ppt_w, ppt_h)
            tf = tx.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
            
            p = tf.paragraphs[0]
            style = item.get('style', {})
            al = style.get('alignment', 'left')
            p.alignment = PP_ALIGN.CENTER if al=='center' else PP_ALIGN.RIGHT if al=='right' else PP_ALIGN.LEFT
            
            run = p.add_run()
            run.text = item.get('content', '')
            run.font.size = Pt(max(8, min(style.get('font_size', 12), 100)))
            run.font.bold = style.get('is_bold', False)
            run.font.color.rgb = hex_to_rgb(style.get('color_hex', '#000000'))
            
        prog.progress((idx+1)/len(image_list))
    
    out = io.BytesIO(); prs.save(out); out.seek(0)
    return out

# ================= 界面 UI =================
st.set_page_config(page_title="Smart Crop PPTX", page_icon="✂️", layout="wide")
st.title("✂️ 智能截图 PPTX (最终版)")
st.caption("支持 PDF / 多图 / 粘贴板 | 自动修复图片边缘")

with st.sidebar:
    key = st.text_input("Gemini API Key", type="password")
    model = st.selectbox("Model", ["gemini-2.5-flash", "gemini-2.0-flash"])

# 布局：左上传，右粘贴
c1, c2 = st.columns(2)
with c1: 
    files = st.file_uploader("拖拽文件 (PDF/PNG)", type=['pdf','png','jpg'], accept_multiple_files=True)
with c2: 
    st.info("或者截图后点击下方:")
    # 粘贴按钮
    paste = paste_image_button("📋 粘贴图片 (Ctrl+V)", background_color="#F0F2F6", errors="ignore")

# 只有当有文件或有粘贴内容时才执行
if (files or paste.image_data) and key:
    st.divider()
    
    # 提取图片
    imgs = extract_images(files, paste.image_data)
    
    if imgs:
        st.write(f"📊 待处理: **{len(imgs)}** 页")
        
        # 预览区
        cols = st.columns(min(5, len(imgs)))
        for i, im in enumerate(imgs[:5]):
            # --- 【这里是修改的地方】 ---
            # 旧写法: use_container_width=True
            # 新写法: width="stretch"
            cols[i].image(im, caption=f"P{i+1}", width="stretch") 
        
        if st.button("🚀 开始转换", type="primary"):
            try:
                data = create_pptx(imgs, key, model)
                st.success("转换完成!")
                st.download_button("📥 下载 PPTX", data, "final_presentation.pptx")
            except Exception as e:
                st.error(f"发生错误: {e}")

elif not key: 
    st.warning("👈 请先在左侧填入 API Key")